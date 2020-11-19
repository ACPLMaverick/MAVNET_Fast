#!/usr/bin/python

import math
import random
import glob
from enum import Enum
import os
import pptx
import pptx.util
import pptx.dml.color
import pptx.enum.text
import tkinter
import tkinter.filedialog
import tkinter.messagebox
import tkinter.colorchooser
from PIL import ImageTk, Image
from tkinter import ttk
import fitz


# Known bugs:
# - Hang when loading images - make it async?

# Add possibility to drag-drop files.
# Make program remember last used paths and parameters

# ###### LOGIC

class enum_enhanced(Enum):
    def __str__(self):
        return self.name.replace("_", " ").capitalize()

    def get_names(enum_class):
        return [str(e) for e in enum_class]

    def reconvert_name(name):
        return name.upper().replace(" ", "_")


class word_orientation(enum_enhanced):
    BOTTOM = 0,
    BOTTOM_LEFT = 1
    BOTTOM_RIGHT = 2,
    TOP_LEFT = 3,
    TOP_RIGHT = 4


class mode_distribution(enum_enhanced):
    NO_WORDS = 0
    ALL_WORDS = 1,
    ONE_WORD_PER_SLIDE = 2,
    ALL_MODES_COMBINED = 3


class mode_work_slide(enum_enhanced):
    NEW_SLIDE = 0,
    SOURCE_SLIDE = 1


class input_data_load_result(Enum):
    FAILURE = 0,
    OK = 1,
    OK_NON_UNIFORM_DIMENSIONS = 2


class var_helper:
    def is_var_valid(var):
        try:
            var.get()
            return True
        except Exception:
            return False


class dimensions(tkinter.StringVar):
    def __init__(self, x, y):
        self.x = tkinter.IntVar(value=x)
        self.x.trace("w", lambda var, index, mode: self.set(self._get_str()))
        self.y = tkinter.IntVar(value=y)
        self.y.trace("w", lambda var, index, mode: self.set(self._get_str()))
        super().__init__(value=self._get_str())

    def is_valid(self):
        return var_helper.is_var_valid(self.x) and var_helper.is_var_valid(self.y)

    def _get_str(self):
        val_x = self.x.get() if var_helper.is_var_valid(self.x) else "???"
        val_y = self.y.get() if var_helper.is_var_valid(self.y) else "???"
        return "{} x {}".format(val_x, val_y)

    def set_dims(self, new_x, new_y):
        self.x.set(new_x)
        self.y.set(new_y)


# Top-left anchor.
class rect:
    def __init__(self, x=0, y=0, width=0, height=0):
        self.x = int(x)
        self.y = int(y)
        self.width = int(width)
        self.height = int(height)
        self.area = width * height

    def create_scaled_with_canvas(in_rect, canvas, reference_width, reference_height):
        scale_x = canvas.winfo_width() / reference_width
        scale_y = canvas.winfo_height() / reference_height

        canv_x = (in_rect.x * scale_x)
        canv_y = (in_rect.y * scale_y)
        canv_w = (in_rect.width * scale_x)
        canv_h = (in_rect.height * scale_y)

        return rect(canv_x, canv_y, canv_w, canv_h)

    def is_zero(self):
        return self.width == 0 or self.height == 0

    def is_intersecting_me(self, other):
        if self.is_zero():
            return False
        self_x2 = self.x + self.width
        self_y2 = self.y + self.height
        other_x2 = other.x + other.width
        other_y2 = other.y + other.height
        return other_x2 >= self.x and other.x <= self_x2 and other_y2 >= self.y and other.y <= self_y2

    def is_inside_me(self, other):
        if self.is_zero():
            return False
        return other.x >= self.x and (other.x + other.width) <= (self.x + self.width) and other.y >= self.y and (other.y + other.height) <= (self.y + self.height)

    def is_exceeding_me_right(self, other):
        return (other.x + other.width) > (self.x + self.width)

    def is_exceeding_me_down(self, other):
        return (other.y + other.height) > (self.y + self.height)

    def __str__(self):
        return "x: {}, y: {}, width: {}, height: {}, area: {}".format(self.x, self.y, self.width, self.height, self.area)


class enum_var(tkinter.StringVar):
    def __init__(self, enum_class, *args, **kwargs):
        self.enum_class = enum_class
        super().__init__(*args, **kwargs)

    def get_enum(self):
        return self.enum_class[self.enum_class.reconvert_name(self.get())]


class distribution_result:
    def __init__(self, grid_size_x, grid_size_y, scaled_dims_x, scaled_dims_y):
        self.grid_size_x = grid_size_x
        self.grid_size_y = grid_size_y
        self.scaled_dims_x = scaled_dims_x
        self.scaled_dims_y = scaled_dims_y


# slide_image's anchor is NW
class slide_image:
    def __init__(self, file_path, idx):
        self._file_path = file_path
        self._img_ori = Image.open(file_path)
        self._img_scaled_init = None
        self._img = None
        self._canvas_img = 0
        self._canvas_text = 0
        self.original_width = self._img_ori.width
        self.original_height = self._img_ori.height
        self.scaled_width = 0
        self.scaled_height = 0
        self.x = 0
        self.y = 0
        self.is_draw_name = False
        self.name = "(no name {})".format(idx)
        self.name_font = None
        self.name_size = 0
        self.name_color = "#000000"

    def compute_text_rect(screen_width, screen_height, text_height, padding_height, text_orientation, mode):
        if mode is not mode_distribution.ONE_WORD_PER_SLIDE:
            return rect()
        height = 2 * padding_height + text_height
        width = screen_width * 0.2
        x = 0
        y = 0
        if text_orientation is word_orientation.BOTTOM:
            x = screen_width // 2 - width // 2
        if text_orientation is word_orientation.TOP_RIGHT or text_orientation is word_orientation.BOTTOM_RIGHT:
            x = screen_width - width
        if text_orientation is word_orientation.BOTTOM or text_orientation is word_orientation.BOTTOM_LEFT or text_orientation is word_orientation.BOTTOM_RIGHT:
            y = screen_height - height
        return rect(x, y, width, height)

    def distribute(images, screen_width, screen_height,
                   image_width, image_height,
                   padding_width, padding_height,
                   mode, text_orientation, text_height, text_font, text_color, is_fill_x):
        # So at this point we have both presentation loaded, slide dimensions available (in points)
        # And also images loaded with their dimensions (in pixels).
        # Scaled image dimensions will be in POINTS, so they can be easily integrated in pptx.

        # TODO HACK Disable uniform X distribution if the mode is one slide per word
        is_fill_x &= mode != mode_distribution.ONE_WORD_PER_SLIDE
        # ~TODO HACK

        num_images = len(images)
        if num_images == 0:
            return False

        screen_rect = rect(0, 0, screen_width, screen_height)
        text_rect = slide_image.compute_text_rect(screen_width, screen_height, text_height, padding_height, text_orientation, mode)
        # text_rect = rect()
        # print("text_rect", text_rect)

        def compute_scale_factor(a_num_rows):
            return ((screen_height - (a_num_rows + 1) * padding_height - (a_num_rows * text_height_per_row)) / image_height) / a_num_rows

        num_rows = 1
        scale_factor = 1.0
        text_height_per_row = text_height if mode == mode_distribution.ALL_WORDS else 0
        while(True):
            # Scale one image to fit exactly num_rows
            scale_factor = compute_scale_factor(num_rows)
            image_scaled_width = image_width * scale_factor
            # Check if all images will fit in this one row. If not, we'll try to fit them in more rows.
            total_images_width = (image_scaled_width + padding_width) * (num_images // num_rows)
            if total_images_width > screen_width:
                num_rows += 1
            else:
                break

        num_cols = math.ceil(num_images / num_rows)

        need_another_pass = False
        pass_num = 1
        for current_num_rows in range(num_rows, num_rows + 2):
            # print("Distribution pass", pass_num, "NumRows:", current_num_rows)
            scale_factor = compute_scale_factor(current_num_rows)
            image_scaled_width = math.floor(image_width * scale_factor)
            image_scaled_height = math.floor(image_height * scale_factor)

            width_per_image = screen_width / num_cols

            # Now distribute the images
            row = 0
            col = 0
            for image_idx in range(num_images):
                img_rect = rect(0, 0, image_scaled_width, image_scaled_height)
                while True:
                    img_rect.x = padding_width
                    if is_fill_x:
                        img_rect.x += col * width_per_image
                    else:
                        img_rect.x += col * (image_scaled_width + padding_width)
                    img_rect.x = int(img_rect.x)
                    img_rect.y = padding_height + row * (image_scaled_height + padding_height + text_height_per_row)

                    col += 1
                    num_cols = max(col - 1, num_cols)
                    if screen_rect.is_exceeding_me_right(img_rect):
                        # Exceeded to the right, so we have to jump to lower row.
                        col = 0
                        row += 1
                        num_cols = max(col - 1, num_cols)
                        current_num_rows = max(row, current_num_rows)
                    elif screen_rect.is_exceeding_me_down(img_rect):
                        # Exceeded down - reached the end of the screen. Distribution failed
                        if need_another_pass is False:
                            need_another_pass = True
                        else:
                            print("Could not distribute image", image_idx, "of rect:", img_rect)
                        break
                    else:
                        # All good? Check if we don't intersect text rect.
                        if not text_rect.is_intersecting_me(img_rect):
                            break
                        # else continue moving img_rect until we fit.
                # if not (img_rect.x >= 0 and (img_rect.x + img_rect.width) < screen_width):
                #     print("Image", image_idx, "exceeds X bounds:", screen_width, img_rect.x, img_rect.width)
                # if not (img_rect.y >= 0 and (img_rect.y + img_rect.height) < screen_height):
                #     print("Image", image_idx, "exceeds Y bounds:", screen_height, img_rect.y, img_rect.height)

                if current_num_rows == num_rows and need_another_pass:
                    break

                # print("Image:", image_idx, "Rect:", img_rect)

                img = images[image_idx]
                img.x = img_rect.x
                img.y = img_rect.y
                img.scaled_width = img_rect.width
                img.scaled_height = img_rect.height
                img.is_draw_name = mode == mode_distribution.ALL_WORDS
                img.name_font = text_font
                img.name_size = text_height
                img.name_color = text_color
            if need_another_pass is False:
                break
            pass_num += 1

        num_rows = current_num_rows
        return distribution_result(num_cols, num_rows, image_scaled_width, image_scaled_height)

    def _resize_img(self, canvas, new_width, new_height):
        img = None
        if self._img_scaled_init is None:
            # This is an optimization to prevent re-scaling the original big image on any parameter change.
            self._img_scaled_init = self._img_ori.resize((new_width, new_height), Image.ANTIALIAS)
            img = self._img_scaled_init
        else:
            img = self._img_scaled_init.resize((new_width, new_height), Image.ANTIALIAS)
        self._img = ImageTk.PhotoImage(img)
        canvas.itemconfig(self._canvas_img, image=self._img)

    def is_under_cursor(self, canvas, mouse_x, mouse_y, reference_width, reference_height):
        canv = rect.create_scaled_with_canvas(rect(self.x, self.y, self.scaled_width, self.scaled_height),
                                              canvas, reference_width, reference_height)
        canv_x2 = canv.x + canv.width
        canv_y2 = canv.y + canv.height
        if self.is_draw_name:
            canv_y2 += int(self.name_size * (canvas.winfo_height() / reference_height))
        return mouse_x > canv.x and mouse_x < canv_x2 and mouse_y > canv.y and mouse_y < canv_y2

    def add_to_slide(self, slide):
        slide.shapes.add_picture(self._file_path,
                                 pptx.util.Pt(self.x),
                                 pptx.util.Pt(self.y),
                                 width=pptx.util.Pt(self.scaled_width),
                                 height=pptx.util.Pt(self.scaled_height))
        if self.is_draw_name:
            text_x = self.x - self.scaled_width // 2
            text_y = self.y + self.scaled_height
            text_w = self.scaled_width * 2          # How to calculate real width?
            text_h = self.name_size * 2

            text_box = slide.shapes.add_textbox(pptx.util.Pt(text_x),
                                                pptx.util.Pt(text_y),
                                                pptx.util.Pt(text_w),
                                                pptx.util.Pt(text_h))
            paragraph = text_box.text_frame.paragraphs[0]
            paragraph.text = self.name
            paragraph.alignment = pptx.enum.text.PP_ALIGN.CENTER
            paragraph.font.size = pptx.util.Pt(int(self.name_size * 1.2))
            paragraph.font.color.rgb = pptx.dml.color.RGBColor.from_string(self.name_color[1:])
            paragraph.font.name = self.name_font

    def init_draw(self, canvas, reference_width, reference_height):
        self._canvas_img = canvas.create_image(0, 0, image=self._img, anchor=tkinter.NW)
        self._canvas_text = canvas.create_text(0, 0, text="", anchor=tkinter.N)
        self.draw(canvas, reference_width, reference_height)

    def draw(self, canvas, reference_width, reference_height):
        canv = rect.create_scaled_with_canvas(rect(self.x, self.y, self.scaled_width, self.scaled_height),
                                              canvas, reference_width, reference_height)

        canvas.coords(self._canvas_img, canv.x, canv.y)

        self._resize_img(canvas, canv.width, canv.height)

        if self.is_draw_name:
            text_size_scale = canvas.winfo_height() / reference_height
            font_tuple = (self.name_font, str(int(self.name_size * text_size_scale)))
            canvas.itemconfigure(self._canvas_text, text=self.name, font=font_tuple, fill=self.name_color)
            text_x = canv.x + (canv.width / 2)
            text_y = canv.y + canv.height
            canvas.coords(self._canvas_text, text_x, text_y)
        else:
            canvas.itemconfigure(self._canvas_text, text="")


class pdf_extractor:
    def __init__(self, file):
        self.file = file
        mtime = os.path.getmtime(self.file)
        self.output_directory = os.path.join(".distribute_images", os.path.splitext(os.path.basename(self.file))[0] + "_{}".format(mtime))

    def obtain_num_images(self):
        num_images = 0
        with fitz.open(self.file) as doc:
            for i in range(len(doc)):
                for img in doc.getPageImageList(i):
                    num_images += 1
        return num_images

    def extract(self, idx_start, idx_end):
        if os.path.isdir(self.output_directory) is False:
            os.makedirs(self.output_directory)
        img_idx = 0
        with fitz.open(self.file) as doc:
            for i in range(len(doc)):
                for img in doc.getPageImageList(i):
                    if img_idx >= idx_start and img_idx <= idx_end:
                        xref = img[0]
                        file_path = os.path.join(self.output_directory, "image_{:03d}_{:04d}.png".format(i, xref))
                        if os.path.isfile(file_path) is False:
                            pix = fitz.Pixmap(doc, xref)
                            if pix.n >= 4:
                                pix = fitz.Pixmap(fitz.csRGB, pix)
                            pix.writePNG(file_path)
                            pix = None
                    img_idx += 1


class distributor:
    def __init__(self):
        self.pres = None
        self.input_file_paths = None
        self.canvas = None

        self.slide_names = ["(no slides)"]

        self.num_images = tkinter.IntVar(value=0)
        self.image_grid_dims = dimensions(0, 0)
        self.image_dims = dimensions(0, 0)
        self.image_scaled_dims = dimensions(0, 0)

        self.source_slide_name = tkinter.StringVar(value=self.slide_names[0])
        self.source_slide_name.trace("w", lambda name, index, mode, self=self: self._on_source_slide_name_changed())
        self.word_font = tkinter.StringVar(value="Impact")
        self.word_size = tkinter.IntVar(value=18)
        self.is_fill_x = tkinter.IntVar(value=1)
        self.word_position = enum_var(enum_class=word_orientation, value=word_orientation.BOTTOM)
        self.operation = enum_var(enum_class=mode_work_slide, value=mode_work_slide.NEW_SLIDE)
        self.mode = enum_var(enum_class=mode_distribution, value=mode_distribution.ALL_MODES_COMBINED)
        self.slide_background_color = tkinter.StringVar(value="#ffffff")
        self.image_padding = dimensions(32, 32)

        self.word_color = tkinter.StringVar(value="#000000")
        self.oneslide_word_size = tkinter.IntVar(value=40)
        self.oneslide_word_padding = dimensions(32, 32)
        self.oneslide_word_active_index = tkinter.IntVar(value=0)

        self.screen_dims = dimensions(0, 0)
        self._pres_file_path = None
        self._source_slide_idx = 0
        self._images = []

        self._canvas_word_id = 0

        self._trace_image_padding = None
        self._trace_mode = None
        self._trace_word_position = None
        self._trace_word_size = None
        self._trace_is_fill_x = None
        self._trace_word_font = None
        self._trace_word_color = None
        self._trace_oneslide_word_size = None
        self._trace_oneslide_word_padding = None
        self._trace_oneslide_word_active_index = None

    def _extract_from_pdf(self):
        print("TODO Implement.")
        return ""

    def _load_pres_data(self):
        self._disable_param_traces()

        if len(self.pres.slides) > 0:
            self.slide_names.clear()
            ctr = 0
            for slide in self.pres.slides:
                name = slide.name
                if name is None or len(name) == 0:
                    self.slide_names.append("Slide {}".format(ctr))
                else:
                    self.slide_names.append(name)
                ctr += 1
            self.source_slide_name.set(self.slide_names[-1])
        else:
            return False

        slide_width = pptx.util.Emu(self.pres.slide_width)
        slide_height = pptx.util.Emu(self.pres.slide_height)
        self.screen_dims.set_dims(int(slide_width.pt), int(slide_height.pt))

        # last_slide = self.pres.slides[-1]
        # if last_slide.background is not None:
        #    if last_slide.background.fill is not None:
        #        if last_slide.background.fill.fore_color is not None:
        #            color = last_slide.background.fill.fore_color
        #            print(color.rgb)
        #            if hasattr(color, "rgb"):
        #                rgb = color.rgb
        #                if rgb is not None:
        #                    print(str(rgb))

        if self.is_all_set():
            if self._setup_images() is False:
                return False

        return True

    def _load_input_data(self):
        self._disable_param_traces()

        self._images.clear()

        self.num_images.set(len(self.input_file_paths))

        if(self.num_images.get() == 0):
            self.image_grid_dims.set_dims(0, 0)
            self.image_dims.set_dims(0, 0)
            self.image_scaled_dims.set_dims(0, 0)
            return input_data_load_result.FAILURE

        return_code = input_data_load_result.OK
        average_width = 0
        average_height = 0
        reference_width = 0
        reference_height = 0
        image_idx = 1
        for image_file in self.input_file_paths:
            image = slide_image(image_file, image_idx)
            self._images.append(image)
            image_width = image.original_width
            image_height = image.original_height
            if reference_width == 0:
                reference_width = image_width
            if reference_height == 0:
                reference_height = image_height
            if image_width != reference_width or image_height != reference_height:
                return_code = input_data_load_result.OK_NON_UNIFORM_DIMENSIONS
            average_width += image_width
            average_height += image_height
            image_idx += 1
        average_width = average_width / self.num_images.get()
        average_height = average_height / self.num_images.get()
        self.image_dims.set_dims(average_width, average_height)

        if self.is_all_set():
            setup_result = self._setup_images()
            if setup_result is False:
                return_code = input_data_load_result.FAILURE

        return return_code

    def _setup_images(self):
        self._disable_param_traces()

        if not self._validate_vars():
            self._enable_param_traces()
            return False

        screen_width = self.screen_dims.x.get()
        screen_height = self.screen_dims.y.get()

        image_width = self.image_dims.x.get()
        image_height = self.image_dims.y.get()

        padding_width = self.image_padding.x.get()
        padding_height = self.image_padding.y.get()

        # TODO Font color.
        result = slide_image.distribute(self._images, screen_width, screen_height, image_width, image_height,
                                        padding_width, padding_height, self.mode.get_enum(), self.word_position.get_enum(),
                                        self.word_size.get(), self.word_font.get(), self.word_color.get(),
                                        self.is_fill_x.get() == 1)

        if result is False:
            self._enable_param_traces()
            return False

        # Otherwise, result is distribution_result
        self.image_grid_dims.set_dims(result.grid_size_x, result.grid_size_y)
        self.image_scaled_dims.set_dims(result.scaled_dims_x, result.scaled_dims_y)

        self._enable_param_traces()

        return True

    def _enable_param_traces(self):
        self._trace_image_padding = self.image_padding.trace("w", lambda name, index, mode, self=self: self._on_edit_traced())
        self._trace_mode = self.mode.trace("w", lambda name, index, mode, self=self: self._on_edit_traced())
        self._trace_word_position = self.word_position.trace("w", lambda name, index, mode, self=self: self._on_edit_traced())
        self._trace_word_size = self.word_size.trace("w", lambda name, index, mode, self=self: self._on_edit_traced())
        self._trace_is_fill_x = self.is_fill_x.trace("w", lambda name, index, mode, self=self: self._on_edit_traced())
        self._trace_word_font = self.word_font.trace("w", lambda name, index, mode, self=self: self._on_edit_traced())
        self._trace_word_color = self.word_color.trace("w", lambda name, index, mode, self=self: self._on_edit_traced())
        self._trace_oneslide_word_size = self.oneslide_word_size.trace("w", lambda name, index, mode, self=self: self._on_edit_traced())
        self._trace_oneslide_word_padding = self.oneslide_word_padding.trace("w", lambda name, index, mode, self=self: self._on_edit_traced())
        self._trace_oneslide_word_active_index = self.oneslide_word_active_index.trace("w", lambda name, index, mode, self=self: self._on_edit_traced())

    def _disable_param_traces(self):
        if self._trace_image_padding is not None:
            self.image_padding.trace_vdelete("w", self._trace_image_padding)
            self._trace_image_padding = None
        if self._trace_mode is not None:
            self.mode.trace_vdelete("w", self._trace_mode)
            self._trace_mode = None
        if self._trace_word_position is not None:
            self.word_position.trace_vdelete("w", self._trace_word_position)
            self._trace_word_position = None
        if self._trace_word_size is not None:
            self.word_size.trace_vdelete("w", self._trace_word_size)
            self._trace_word_size = None
        if self._trace_is_fill_x is not None:
            self.is_fill_x.trace_vdelete("w", self._trace_is_fill_x)
            self._trace_is_fill_x = None
        if self._trace_word_font is not None:
            self.word_font.trace_vdelete("w", self._trace_word_font)
            self._trace_word_font = None
        if self._trace_word_color is not None:
            self.word_color.trace_vdelete("w", self._trace_word_color)
            self._trace_word_color = None
        if self._trace_oneslide_word_size is not None:
            self.oneslide_word_size.trace_vdelete("w", self._trace_oneslide_word_size)
            self._trace_oneslide_word_size = None
        if self._trace_oneslide_word_padding is not None:
            self.oneslide_word_padding.trace_vdelete("w", self._trace_oneslide_word_padding)
            self._trace_oneslide_word_padding = None
        if self._trace_oneslide_word_active_index is not None:
            self.oneslide_word_active_index.trace_vdelete("w", self._trace_oneslide_word_active_index)
            self._trace_oneslide_word_active_index = None

    def _on_edit_traced(self):
        if self._setup_images():
            self.draw(self.canvas)

    def _validate_vars(self):
        if not self.screen_dims.is_valid():
            return False
        if not self.image_dims.is_valid():
            return False
        if not self.image_padding.is_valid():
            return False
        if not var_helper.is_var_valid(self.word_size):
            return False
        if not var_helper.is_var_valid(self.word_position):
            return False
        if not var_helper.is_var_valid(self.word_font):
            return False
        if not var_helper.is_var_valid(self.word_color):
            return False
        if not var_helper.is_var_valid(self.oneslide_word_size):
            return False
        if not self.oneslide_word_padding.is_valid():
            return False
        if not var_helper.is_var_valid(self.oneslide_word_active_index):
            return False
        return True

    def _create_shuffled_image_indices(self):
        num_indices = len(self._images)
        randomized = list(range(num_indices))
        random.shuffle(randomized)
        return randomized

    def _create_new_slide(self):
        source_slide_layout = self.pres.slides[self._source_slide_idx].slide_layout
        return self.pres.slides.add_slide(source_slide_layout)

    def _acquire_first_slide(self):
        op = self.operation.get_enum()
        if op == mode_work_slide.NEW_SLIDE:
            return self._create_new_slide()
        elif op == mode_work_slide.SOURCE_SLIDE:
            slide = self.pres.slides[self._source_slide_idx]
            # slide.shapes.clear()          # Not supported.
            # slide.placeholders.clear()    # Not supported.
            return slide
        else:
            print("Error:", "Unknown mode operation.")
            return None

    def _add_word_to_slide(self, slide, word_idx):
        screen_width = self.screen_dims.x.get()
        screen_height = self.screen_dims.y.get()
        word_size = self.oneslide_word_size.get()
        word_padding_x = self.oneslide_word_padding.x.get()
        word_padding_y = self.oneslide_word_padding.y.get()
        word_alignment = pptx.enum.text.PP_ALIGN.LEFT
        word_x = word_padding_x
        word_y = word_padding_y
        word_w = screen_width // 2
        word_h = word_size * 2

        if self.word_position.get_enum() == word_orientation.TOP_RIGHT:
            word_alignment = pptx.enum.text.PP_ALIGN.RIGHT
            word_x = screen_width - word_padding_x - word_w
            word_y = word_padding_y
        elif self.word_position.get_enum() == word_orientation.BOTTOM_RIGHT:
            word_alignment = pptx.enum.text.PP_ALIGN.RIGHT
            word_x = screen_width - word_padding_x - word_w
            word_y = screen_height - word_padding_y - word_h
        elif self.word_position.get_enum() == word_orientation.BOTTOM_LEFT:
            word_alignment = pptx.enum.text.PP_ALIGN.LEFT
            word_x = word_padding_x
            word_y = screen_height - word_padding_y - word_h
        elif self.word_position.get_enum() == word_orientation.TOP_LEFT:
            word_alignment = pptx.enum.text.PP_ALIGN.LEFT
            word_x = word_padding_x
            word_y = word_padding_y
        elif self.word_position.get_enum() == word_orientation.BOTTOM:
            word_alignment = pptx.enum.text.PP_ALIGN.CENTER
            word_x = screen_width // 2 - word_w // 2
            word_y = screen_height - word_padding_y - word_h
        else:
            print("Not supported word orientation.")

        text_box = slide.shapes.add_textbox(pptx.util.Pt(word_x),
                                            pptx.util.Pt(word_y),
                                            pptx.util.Pt(word_w),
                                            pptx.util.Pt(word_h))
        paragraph = text_box.text_frame.paragraphs[0]
        paragraph.text = self._images[word_idx].name
        paragraph.alignment = word_alignment
        paragraph.font.size = pptx.util.Pt(int(word_size * 1.2))
        paragraph.font.color.rgb = pptx.dml.color.RGBColor.from_string(self.word_color.get()[1:])
        paragraph.font.name = self.word_font.get()

    def _create_all(self, first_slide, use_background):
        last_mode = self.mode.get_enum()

        self.mode.set(mode_distribution.NO_WORDS)   # This should trigger the setup_images call.
        result = self._create_one_mode(first_slide, self.mode.get_enum(), use_background)
        if result is False:
            self.mode.set(last_mode)
            return False

        next_slide = self._create_new_slide()
        self.mode.set(mode_distribution.ONE_WORD_PER_SLIDE)
        result = self._create_one_mode(next_slide, self.mode.get_enum(), use_background)
        if result is False:
            self.mode.set(last_mode)
            return False

        next_slide = self._create_new_slide()
        self.mode.set(mode_distribution.ALL_WORDS)
        result = self._create_one_mode(next_slide, self.mode.get_enum(), use_background)
        if result is False:
            self.mode.set(last_mode)
            return False

        self.mode.set(last_mode)
        return True

    def _create_one_mode(self, slide, mode, use_background):
        if slide is None:
            return False

        for img in self._images:
            img.add_to_slide(slide)

        if use_background:
            rgb_color = pptx.dml.color.RGBColor.from_string(self.slide_background_color.get()[1:])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = rgb_color

        if mode == mode_distribution.ONE_WORD_PER_SLIDE:
            indices = self._create_shuffled_image_indices()
            self._add_word_to_slide(slide, indices[0])
            del indices[0]
            for word_idx in indices:
                new_slide = self._create_new_slide()
                for img in self._images:
                    img.add_to_slide(new_slide)
                if use_background:
                    new_slide.background.fill.solid()
                    new_slide.background.fill.fore_color.rgb = rgb_color
                self._add_word_to_slide(new_slide, word_idx)

        return True

    def _on_source_slide_name_changed(self):
        # Update the idx.
        slide_name = self.source_slide_name.get()
        try:
            self._source_slide_idx = self.slide_names.index(slide_name)
        except ValueError as e:
            print(e)
            print("Not found slide name. Not setting the slide idx.")

    def get_images(self):
        return self._images

    def get_image_under_cursor(self, canvas, mouse_x, mouse_y):
        for img in self._images:
            if img.is_under_cursor(canvas, mouse_x, mouse_y, self.screen_dims.x.get(), self.screen_dims.y.get()):
                return img
        return None

    def is_all_set(self):
        return self.pres is not None and self.input_file_paths is not None and len(self._images) > 0

    def set_pres(self, path):
        if path.lower().endswith(".pptx") or path.lower().endswith(".pptm"):
            self._pres_file_path = path
            self.pres = pptx.Presentation(path)
            return self._load_pres_data()
        else:
            return False

    def set_input(self, input_file_paths):
        for file_path in input_file_paths:
            if os.path.isfile(file_path) is False:
                return input_data_load_result.FAILURE
        self.input_file_paths = input_file_paths
        return self._load_input_data()

    def create(self, use_background):
        if self.pres is None:
            return False

        if self.input_file_paths is None:
            return False

        if self._validate_vars() is False:
            return False

        if len(self._images) <= 0:
            return False

        result = False
        if self.mode.get_enum() == mode_distribution.ALL_MODES_COMBINED:
            result = self._create_all(self._acquire_first_slide(), use_background)
        else:
            result = self._create_one_mode(self._acquire_first_slide(), self.mode.get_enum(), use_background)

        if result:
            self.pres.save(self._pres_file_path)
            # Reload the presentation from disk, to update all visible data.
            self.set_pres(self._pres_file_path)
        return result

    def _draw_word(self, canvas):
        if self.mode.get_enum() == mode_distribution.ONE_WORD_PER_SLIDE:
            text_size_scale = canvas.winfo_height() / self.screen_dims.y.get()
            font_tuple = (self.word_font.get(), str(int(self.oneslide_word_size.get() * text_size_scale)))
            word_padding_x = self.oneslide_word_padding.x.get()
            word_padding_y = self.oneslide_word_padding.y.get()
            word_anchor = tkinter.NW
            word_x = word_padding_x
            word_y = word_padding_y
            if self.word_position.get_enum() == word_orientation.TOP_RIGHT:
                word_anchor = tkinter.NE
                word_x = self.screen_dims.x.get() - word_padding_x
                word_y = word_padding_y
            elif self.word_position.get_enum() == word_orientation.BOTTOM_RIGHT:
                word_anchor = tkinter.SE
                word_x = self.screen_dims.x.get() - word_padding_x
                word_y = self.screen_dims.y.get() - word_padding_y
            elif self.word_position.get_enum() == word_orientation.BOTTOM_LEFT:
                word_anchor = tkinter.SW
                word_x = word_padding_x
                word_y = self.screen_dims.y.get() - word_padding_y
            elif self.word_position.get_enum() == word_orientation.TOP_LEFT:
                word_anchor = tkinter.NW
                word_x = word_padding_x
                word_y = word_padding_y
            elif self.word_position.get_enum() == word_orientation.BOTTOM:
                word_anchor = tkinter.S
                word_x = self.screen_dims.x.get() // 2
                word_y = self.screen_dims.y.get() - word_padding_y
            else:
                print("Not supported word orientation.")
            word_rect = rect.create_scaled_with_canvas(rect(word_x, word_y, 1, 1),
                                                       canvas, self.screen_dims.x.get(), self.screen_dims.y.get())
            image_idx = int(self.oneslide_word_active_index.get()) if var_helper.is_var_valid(self.oneslide_word_active_index) else 1
            canvas.itemconfigure(self._canvas_word_id, text=self._images[image_idx - 1].name, font=font_tuple,
                                 fill=self.word_color.get(), anchor=word_anchor)
            canvas.coords(self._canvas_word_id, word_rect.x, word_rect.y)
        else:
            canvas.itemconfigure(self._canvas_word_id, text="")

    def init_draw(self, canvas):
        self.canvas = canvas
        for img in self._images:
            img.init_draw(canvas, self.screen_dims.x.get(), self.screen_dims.y.get())
        self._canvas_word_id = canvas.create_text(0, 0, text="", anchor=tkinter.NW)

    def draw(self, canvas):
        if canvas is None:
            return
        for img in self._images:
            img.draw(canvas, self.screen_dims.x.get(), self.screen_dims.y.get())
        self._draw_word(canvas)


# ####### UI

# https://stackoverflow.com/questions/51902451/how-to-enable-and-disable-frame-instead-of-individual-widgets-in-python-tkinter/52152773
class tk_util:
    def create_num_input(root, var_value, min_val, max_val):
        def validate_num_input(val, widget):
            int_val = 0
            try:
                int_val = int(val)
            except Exception:
                try:
                    is_zero_length_string = len(val) == 0
                    return is_zero_length_string
                except Exception:
                    return False
                return False
            return int_val >= widget["from"] and int_val <= widget["to"]
        sb = tkinter.Spinbox(root, textvariable=var_value, width=5,
                             from_=min_val, to=max_val, validate="all")
        cmd = root.register(lambda val: validate_num_input(val, sb))
        sb.configure(validatecommand=(cmd, "%P"))
        return sb

    def place_in_grid(widget, x, y, w=1, h=1, weight_w=1, weight_h=0, orientation="", nopadding=False):
        root = widget._root()
        padding = 3 if nopadding is False else 0
        widget.grid(row=y, column=x, rowspan=h, columnspan=w, padx=padding, pady=padding, sticky=orientation)
        root.columnconfigure(x, weight=weight_w)
        root.rowconfigure(y, weight=weight_h)

    def window_set_on_cursor(window):
        rt = window._root()
        x = rt.winfo_pointerx()
        y = rt.winfo_pointery()
        window.geometry("+{}+{}".format(x, y))

    def window_set_on_root(window, orientation=tkinter.CENTER, offset_x=0, offset_y=0):
        rt = window._root()
        x = int(rt.winfo_x()) + offset_x
        y = int(rt.winfo_y()) + offset_y
        if orientation == tkinter.CENTER:
            x += int(rt.winfo_width() / 2)
            y += int(rt.winfo_height() / 2)
        window.geometry("+{}+{}".format(x, y))

    def bind_enter_esc(window, func_enter=None, func_esc=None):
        if func_enter is not None:
            window.bind("<Return>", lambda event: func_enter())
        if func_esc is not None:
            window.bind("<Escape>", lambda event: func_esc())


class labelframe_disableable(tkinter.LabelFrame):
    def enable(self, status="normal"):
        def cstate(widget):
            # Is this widget a container?
            if widget.winfo_children:
                # It's a container, so iterate through its children
                for w in widget.winfo_children():
                    if len(w.winfo_children()) == 0:
                        w.configure(state=status)
                    else:
                        cstate(w)
        cstate(self)

    def disable(self):
        self.enable(status="disabled")


class color_picker(tkinter.Frame):
    def __init__(self, master, var_color, bg_color, width=144, height=21, command=None, can_be_disabled=True):
        super().__init__(master, width=width, height=height, bg=bg_color)
        self._can_be_disabled = can_be_disabled
        width_btn = int(0.8 * width) if can_be_disabled else width
        width_cb = width - width_btn
        if can_be_disabled:
            self._var_check = tkinter.IntVar(value=0)
            self.cb = tkinter.Checkbutton(self, variable=self._var_check, width=width_cb, height=height,
                                          bg=bg_color, highlightcolor=bg_color, activebackground=bg_color)
        self.btn = tkinter.Button(self, command=self._on_clicked, bg=var_color.get(), fg=var_color.get())
        if can_be_disabled:
            self.cb.place(x=0, y=0, width=width_cb, height=height)
        self.btn.place(x=width_cb, y=0, width=width_btn, height=height)
        self.var_color = var_color
        self.var_color.trace("w", lambda name, index, mode, self=self: self._on_color_changed())
        if can_be_disabled:
            self._var_check.trace("w", lambda name, index, mode, self=self: self._on_checkbox_changed())
        self.command = command
        self._var_color_init_val = self.var_color.get()

    def is_enabled(self):
        if self._can_be_disabled:
            return self._var_check.get() == 1
        else:
            return True

    def _on_clicked(self):
        new_color_tuple = tkinter.colorchooser.askcolor(self.var_color.get(), title="Pick a color")
        if len(new_color_tuple) > 1:
            new_color = new_color_tuple[1]
            if new_color is not None and new_color != self.var_color.get():
                self.var_color.set(new_color)
                if self.command is not None:
                    self.command(new_color)
                if self._can_be_disabled:
                    if self._var_check.get() == 0:
                        self._var_check.set(1)

    def _on_color_changed(self):
        self.btn["bg"] = self.var_color.get()
        self.btn["fg"] = self.var_color.get()

    def _on_checkbox_changed(self):
        if self.is_enabled() is False:
            # Reset to default.
            self.var_color.set(self._var_color_init_val)


class panel_edit_single_word(tkinter.Toplevel):
    def __init__(self, color_bg, text, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.color_bg = color_bg
        self.text = text
        self.title = "Edit a word"
        self.is_ok = False
        self.resizable(False, False)
        self.configure(bg=self.color_bg)
        tk_util.bind_enter_esc(self, self._on_ok_clicked, self._on_cancel_clicked)

        self._text_var = tkinter.StringVar(value=self.text)

        self._edit = tkinter.Entry(self, width=20, textvariable=self._text_var)
        self._edit.focus_set()
        self._edit.selection_range(0, tkinter.END)
        self._edit.pack(padx=3, pady=3)

        frame_btns = tkinter.Frame(self, bg=self.color_bg)
        frame_btns.pack()
        self.btn_ok = tkinter.Button(frame_btns, text="OK", width=8, command=self._on_ok_clicked)
        self.btn_ok.pack(side=tkinter.LEFT, padx=3, pady=3)
        self.btn_cancel = tkinter.Button(frame_btns, text="Cancel", width=8, command=self._on_cancel_clicked)
        self.btn_cancel.pack(side=tkinter.LEFT, padx=3, pady=3)

        tk_util.window_set_on_cursor(self)

    def _on_ok_clicked(self):
        self.text = self._text_var.get()
        self.is_ok = True
        self.destroy()

    def _on_cancel_clicked(self):
        self.destroy()


class panel_word_settings(tkinter.Toplevel):
    def __init__(self, color_bg, objects, is_appear_under_cursor=False, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.objects = objects
        self.edited_names = dict()
        self.title = "Word settings"
        self.color_bg = color_bg
        self.is_ok = False
        self.resizable(False, False)
        self.configure(bg=color_bg)
        self._create_widgets()
        tk_util.bind_enter_esc(self, self._on_ok_clicked, self._on_cancel_clicked)
        if is_appear_under_cursor:
            tk_util.window_set_on_cursor(self)
        else:
            tk_util.window_set_on_root(self)

    def _create_widgets(self):
        self.btn_edit = tkinter.Button(self, text="Edit", width=8, command=self._on_edit_clicked)
        tk_util.place_in_grid(self.btn_edit, 0, 0, orientation=tkinter.NW)
        self.list_box = tkinter.Listbox(self, selectmode=tkinter.SINGLE, height=len(self.objects))
        self.list_box.bind("<Double-Button-1>", lambda event, self=self: self._on_edit_clicked())
        tk_util.place_in_grid(self.list_box, 1, 0, orientation=tkinter.NW)

        for obj in self.objects:
            word = obj.name
            self.list_box.insert(tkinter.END, word)

        frame_btns = tkinter.Frame(self, bg=self.color_bg)
        tk_util.place_in_grid(frame_btns, 0, 2, w=2, nopadding=True)
        self.btn_ok = tkinter.Button(frame_btns, text="OK", width=8, command=self._on_ok_clicked)
        self.btn_ok.pack(side=tkinter.LEFT, padx=3, pady=3)
        self.btn_cancel = tkinter.Button(frame_btns, text="Cancel", width=8, command=self._on_cancel_clicked)
        self.btn_cancel.pack(side=tkinter.LEFT, padx=3, pady=3)

    def _on_edit_clicked(self):
        if len(self.list_box.curselection()) < 1:
            return
        selected_idx = self.list_box.curselection()[0]
        text = self.edited_names[selected_idx] if selected_idx in self.edited_names else self.objects[selected_idx].name
        edit_panel = panel_edit_single_word(self.color_bg, text, self)
        edit_panel.grab_set()
        self.wait_window(edit_panel)
        if edit_panel.is_ok:
            self.edited_names[selected_idx] = edit_panel.text
            self.list_box.delete(selected_idx)
            self.list_box.insert(selected_idx, edit_panel.text)

    def _on_ok_clicked(self):
        for key, value in self.edited_names.items():
            self.objects[key].name = value
        self.is_ok = True
        self.destroy()

    def _on_cancel_clicked(self):
        self.destroy()


class panel_extract_pdf_options(tkinter.Toplevel):
    def __init__(self, color_bg, num_images, var_from, var_to, *args, **kwargs):
        super().__init__(*args, *kwargs)
        self.title = "Extract PDF"
        self.is_ok = False
        self.label = tkinter.Label(self,
                                   text="There are {} images in this PDF.\nPlease provide an image range for PDF extraction.".format(num_images),
                                   bg=color_bg)
        tk_util.place_in_grid(self.label, 0, 0, w=4)

        var_from.set(1)
        var_to.set(num_images)

        self.label_from = tkinter.Label(self, text="From:", bg=color_bg)
        tk_util.place_in_grid(self.label_from, 0, 1, orientation=tkinter.E)
        self.edit_from = tk_util.create_num_input(self, var_from, 1, num_images)
        tk_util.place_in_grid(self.edit_from, 1, 1, orientation=tkinter.W)
        self.label_to = tkinter.Label(self, text="To:", bg=color_bg)
        tk_util.place_in_grid(self.label_to, 2, 1, orientation=tkinter.E)
        self.edit_to = tk_util.create_num_input(self, var_to, 1, num_images)
        tk_util.place_in_grid(self.edit_to, 3, 1, orientation=tkinter.W)

        self.btn_ok = tkinter.Button(self, text="OK", width=8, command=self._on_ok_clicked)
        tk_util.place_in_grid(self.btn_ok, 0, 2, orientation=tkinter.E, w=2)
        self.btn_cancel = tkinter.Button(self, text="Cancel", width=8, command=self._on_cancel_clicked)
        tk_util.place_in_grid(self.btn_cancel, 2, 2, orientation=tkinter.W, w=2)

        self.resizable(False, False)
        self.configure(bg=color_bg)
        tk_util.window_set_on_cursor(self)
        tk_util.bind_enter_esc(self, self._on_ok_clicked, self._on_cancel_clicked)

    def _on_ok_clicked(self):
        self.is_ok = True
        self.destroy()

    def _on_cancel_clicked(self):
        self.destroy()


class panel_pick_images(tkinter.Toplevel):
    def __init__(self, color_bg, directory, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._all_image_files = None
        self.selected_image_files = None
        self.title = "Select images"
        self.is_ok = False
        self._gather_image_files(directory)
        self._create_widgets(color_bg, directory)
        self.resizable(False, False)
        self.configure(bg=color_bg)
        tk_util.window_set_on_root(self, orientation=tkinter.NW, offset_x=32, offset_y=32)
        tk_util.bind_enter_esc(self, self._on_ok_clicked, self._on_cancel_clicked)

    def _gather_image_files(self, directory):
        self._all_image_files = glob.glob(os.path.join(directory, "*.jpg"))
        self._all_image_files.extend(glob.glob(os.path.join(directory, "*.png")))
        self._all_image_files.extend(glob.glob(os.path.join(directory, "*.svg")))
        self._all_image_files.extend(glob.glob(os.path.join(directory, "*.tga")))
        self._all_image_files.extend(glob.glob(os.path.join(directory, "*.gif")))

    def _create_widgets(self, color_bg, directory):
        num_image_files = len(self._all_image_files)
        is_have_any_images = num_image_files > 0
        if is_have_any_images:
            self.list_box = tkinter.Listbox(self, selectmode=tkinter.EXTENDED, height=num_image_files)
            tk_util.place_in_grid(self.list_box, 0, 0, w=2)
            width = 0
            for img_file in self._all_image_files:
                path_in_listbox = os.path.relpath(img_file, directory)
                path_in_listbox_num_chars = len(path_in_listbox)
                if path_in_listbox_num_chars > width:
                    width = path_in_listbox_num_chars
                self.list_box.insert(tkinter.END, path_in_listbox)
            width = max(width, 40)
            self.list_box.configure(width=width)
            self.list_box.bind("<<ListboxSelect>>", lambda event, self=self: self._on_selection_changed())
            self.list_box.selection_set(0, tkinter.END)

            self.btn_ok = tkinter.Button(self, text="OK", width=8, command=self._on_ok_clicked)
            tk_util.place_in_grid(self.btn_ok, 0, 1, orientation=tkinter.E)
            self.btn_cancel = tkinter.Button(self, text="Cancel", width=8, command=self._on_cancel_clicked)
            tk_util.place_in_grid(self.btn_cancel, 1, 1, orientation=tkinter.W)
        else:
            self.label = tkinter.Label(self, text="There are no images in this directory.\nPlease select a different one.", bg=color_bg)
            tk_util.place_in_grid(self.label, 0, 0)
            self.btn_ok = tkinter.Button(self, text="OK", width=8, command=self._on_cancel_clicked)
            tk_util.place_in_grid(self.btn_ok, 0, 1)

    def _on_selection_changed(self):
        curr_selection_len = len(self.list_box.curselection())
        btn_ok_state = self.btn_ok["state"]
        if curr_selection_len == 0 and btn_ok_state != "disabled":
            self.btn_ok.configure(state="disabled")
        elif curr_selection_len != 0 and btn_ok_state == "disabled":
            self.btn_ok.configure(state="normal")

    def _on_ok_clicked(self):
        selection = self.list_box.curselection()
        self.selected_image_files = []
        for index in selection:
            self.selected_image_files.append(self._all_image_files[index])
        self.is_ok = True
        self.destroy()

    def _on_cancel_clicked(self):
        self.destroy()


class window:
    _color_bg = "#fafafa"
    _str_select_pres = "Please select a presentation file."
    _str_select_input = "Please select a PDF file or an image directory."

    def __init__(self):
        self._create_top()
        self.distrib = distributor()
        self._create_widgets()
        self._lock_edit()

    def _create_top(self):
        self.width = 770
        self.height = 705

        self.top = tkinter.Tk()
        self.top.title("Distribute Images Tool")
        self.top.resizable(False, False)
        self.top.geometry("{}x{}".format(self.width, self.height))
        self.top.configure(bg=window._color_bg)
        # self.top.columnconfigure(0, minsize=width)
        tk_util.bind_enter_esc(self.top, func_esc=self.conditional_exit)
        self.top.protocol("WM_DELETE_WINDOW", self.conditional_exit)

    def _create_widgets(self):
        # Directory with images (TODO: PDF file for automatic extraction)

        # After load:
        # # Either update selected slide or create a new one based on selected slide
        # # Total number of images and number n x n
        # # Spacing size between them (recalculates size and n x n)
        # # Recalculated (display warning when they are different and take average)
        # # Word positioning
        # # Word font
        # # Slide background
        # # Preview window? Display preview using simple rectangle shapes
        # # Create button

        #####################

        # File picker

        frame_pick_files = labelframe_disableable(self.top, bg=window._color_bg)
        tk_util.place_in_grid(frame_pick_files, 0, 0, w=2)

        self.label_pres = tkinter.Label(frame_pick_files, width=92, text=window._str_select_pres,
                                        anchor=tkinter.W, bg=window._color_bg)
        tk_util.place_in_grid(self.label_pres, 0, 0)

        self.btn_select_pres = tkinter.Button(frame_pick_files, width=12, text="Presentation", command=self._cmd_select_file_pptx)
        tk_util.place_in_grid(self.btn_select_pres, 1, 0, orientation=tkinter.E)

        self.label_dir = tkinter.Label(frame_pick_files, width=92, text=window._str_select_input,
                                       anchor=tkinter.W, bg=window._color_bg)
        tk_util.place_in_grid(self.label_dir, 0, 1)

        frame_two_button = tkinter.Frame(frame_pick_files, bg=window._color_bg)
        tk_util.place_in_grid(frame_two_button, 1, 1, orientation=tkinter.E)

        self.btn_select_file = tkinter.Button(frame_two_button, text="PDF", width=5, command=self._cmd_select_file_pdf)
        tk_util.place_in_grid(self.btn_select_file, 0, 0, orientation=tkinter.W, nopadding=True)

        self.btn_select_dir = tkinter.Button(frame_two_button, text="Images", command=self._cmd_select_dir)
        tk_util.place_in_grid(self.btn_select_dir, 1, 0, orientation=tkinter.E, nopadding=True)

        #####################

        # Visual editor

        canvas_size_mul = 0.95
        self.canvas = tkinter.Canvas(self.top, bg=self.distrib.slide_background_color.get(),
                                     width=self.width * canvas_size_mul, height=(self.width * (9 / 16)) * canvas_size_mul)
        tk_util.place_in_grid(self.canvas, 0, 1, w=2)
        self._configure_canvas()

        #####################

        # Editing

        self.frame_edit = labelframe_disableable(self.top, bg=window._color_bg)
        tk_util.place_in_grid(self.frame_edit, 0, 2)

        # Distribution mode
        # Source slide
        # Screen size
        # Image padding
        # Word positioning
        # Word font and size (latter adjusts automatically to distrib. mode)

        self.edit_source_slide = window._w_create_pair_combobox(self.frame_edit, "Source slide:", [self.distrib.source_slide_name.get()],
                                                                self.distrib.source_slide_name, 0, 0)
        # TODO Add more font support. For now let's support just one.
        self.edit_word_font = window._w_create_pair_combobox(self.frame_edit, "Word font:", [self.distrib.word_font.get()],
                                                             self.distrib.word_font, 0, 1)
        frame_size_and_fill_x = tkinter.Frame(self.frame_edit, bg=window._color_bg)
        self.edit_word_font_size = window._w_create_pair_num_edit(frame_size_and_fill_x, "Small text:", self.distrib.word_size, 0, 0, 1, 64)
        self.edit_word_large_size = window._w_create_pair_num_edit(frame_size_and_fill_x, "Large text:", self.distrib.oneslide_word_size, 2, 0, 1, 110)
        tk_util.place_in_grid(frame_size_and_fill_x, 0, 2, w=2, orientation=tkinter.W, nopadding=True)

        self.edit_word_positioning = window._w_create_pair_combobox(self.frame_edit, "Word position:",
                                                                    enum_enhanced.get_names(word_orientation), self.distrib.word_position,
                                                                    0, 3)
        self.edit_operation = window._w_create_pair_combobox(self.frame_edit, "Work slide:", enum_enhanced.get_names(mode_work_slide),
                                                             self.distrib.operation, 2, 0)
        self.edit_mode = window._w_create_pair_combobox(self.frame_edit, "Distribution:", enum_enhanced.get_names(mode_distribution),
                                                        self.distrib.mode, 2, 1)
        self.edit_bg_color = window._w_create_pair_color_picker(self.frame_edit, "Background color:", self.distrib.slide_background_color, 2, 2)
        self.edit_font_color = window._w_create_pair_color_picker(self.frame_edit, "Font color:", self.distrib.word_color, 2, 3, False)

        self.edit_is_fill_x = window._w_create_pair_checkbox(self.frame_edit, "Uniform width:", self.distrib.is_fill_x, 0, 4)
        self.edit_padding = window._w_create_pair_dim_edit(self.frame_edit, "Image padding (pt):", self.distrib.image_padding, 2, 4, 1, 128)

        self.edit_active_word = window._w_create_pair_num_edit(self.frame_edit, "Active word:",
                                                               self.distrib.oneslide_word_active_index, 0, 5)
        self.edit_word_padding = window._w_create_pair_dim_edit(self.frame_edit, "Word padding (pt):", self.distrib.oneslide_word_padding,
                                                                2, 5, 1, 128)

        # #####################

        # Informations

        self.frame_info = labelframe_disableable(self.top, bg=window._color_bg)
        tk_util.place_in_grid(self.frame_info, 1, 2)

        self.label_number = window._w_create_pair_label(self.frame_info, "Number of images:", self.distrib.num_images, 0, 0)
        self.label_number_nxn = window._w_create_pair_label(self.frame_info, "Grid dimensions (WxH):", self.distrib.image_grid_dims, 0, 1)
        self.label_dimensions_screen = window._w_create_pair_label(self.frame_info, "Screen dimensions (pt):", self.distrib.screen_dims, 0, 2)
        self.label_dimensions_ori = window._w_create_pair_label(self.frame_info, "Ori. image dimensions (px):", self.distrib.image_dims, 0, 3)
        self.label_dimensions_rec = window._w_create_pair_label(self.frame_info, "Scl. image dimensions (pt):", self.distrib.image_scaled_dims, 0, 4)

        self.btn_word_settings = tkinter.Button(self.frame_info, text="Word settings", width=25, command=lambda: self._canv_open_edit_words(False))
        tk_util.place_in_grid(self.btn_word_settings, 0, 5, w=2)

        #####################

        # Create button at the very end.

        self.btn_create = tkinter.Button(self.top, text="Create", width=34, command=self._cmd_create)
        tk_util.place_in_grid(self.btn_create, 0, 3, 2)

    def _lock_edit(self):
        self.frame_edit.disable()
        self.frame_info.disable()
        self.btn_create.configure(state="disabled")

    def _unlock_edit(self):
        self.frame_edit.enable()
        self.frame_info.enable()
        self.btn_create.configure(state="normal")

    def _all_set(self):
        self.edit_active_word.configure(to=len(self.distrib.get_images()))
        self._unlock_edit()
        self.distrib.init_draw(self.canvas)

    def _configure_canvas(self):
        self.canvas.bind("<Double-Button-1>", self._canv_on_left_double_clicked)
        self.distrib.slide_background_color.trace("w", lambda name, index, mode, self=self: self._canv_on_slide_background_color_changed())

        # Right-click menu
        m = tkinter.Menu(self.canvas, tearoff=0)
        m.add_command(label="Word settings", command=lambda: self._canv_open_edit_words(True))

        def do_popup(event):
            if self.distrib.is_all_set() is False:
                return
            try:
                m.tk_popup(event.x_root, event.y_root)
            finally:
                m.grab_release()
        self.canvas.bind("<Button-3>", do_popup)

    def _canv_on_left_double_clicked(self, event):
        if self.distrib.is_all_set() is False:
            return
        slide_clicked = self.distrib.get_image_under_cursor(self.canvas, event.x, event.y)
        if slide_clicked is None:
            return
        edit_panel = panel_edit_single_word(self._color_bg, slide_clicked.name, self.top)
        edit_panel.grab_set()
        self.top.wait_window(edit_panel)
        if edit_panel.is_ok:
            slide_clicked.name = edit_panel.text
            self.distrib.draw(self.canvas)

    def _canv_open_edit_words(self, is_appear_under_cursor):
        objects = self.distrib.get_images()
        edit_words = panel_word_settings(window._color_bg, objects, is_appear_under_cursor, self.top)
        edit_words.grab_set()
        self.top.wait_window(edit_words)
        if edit_words.is_ok:
            self.distrib._setup_images()
            self.distrib.draw(self.canvas)

    def _canv_on_slide_background_color_changed(self):
        self.canvas["bg"] = self.distrib.slide_background_color.get()

    def _w_create_pair(root, element, text_name, base_x, base_y):
        label_name = tkinter.Label(root, text=text_name, bg=window._color_bg)
        tk_util.place_in_grid(label_name, base_x, base_y, orientation=tkinter.W)
        tk_util.place_in_grid(element, base_x + 1, base_y, orientation=tkinter.E)
        return element

    def _w_create_pair_label(root, text_name, var_value, base_x, base_y):
        return window._w_create_pair(root, tkinter.Label(root, textvariable=var_value, bg=window._color_bg),
                                     text_name, base_x, base_y)

    def _w_create_pair_num_edit(root, text_name, var_value, base_x, base_y, min_val=1, max_val=9999):
        dim = tk_util.create_num_input(root, var_value, min_val, max_val)
        return window._w_create_pair(root, dim, text_name, base_x, base_y)

    def _w_create_pair_dim_edit(root, text_name, var_dim, base_x, base_y, min_val=1, max_val=9999):
        inter_frame = tkinter.Frame(root, bg=window._color_bg)
        dim_x = tk_util.create_num_input(inter_frame, var_dim.x, min_val, max_val)
        dim_y = tk_util.create_num_input(inter_frame, var_dim.y, min_val, max_val)
        sep = tkinter.Label(inter_frame, text="x", bg=window._color_bg)
        tk_util.place_in_grid(dim_x, 0, 0, nopadding=True)
        tk_util.place_in_grid(sep, 1, 0, nopadding=True)
        tk_util.place_in_grid(dim_y, 2, 0, nopadding=True)
        return window._w_create_pair(root, inter_frame, text_name, base_x, base_y)

    def _w_create_pair_combobox(root, text_name, options, var_value, base_x, base_y):
        def func_validate():
            return False
        cmd = root.register(func_validate)
        combo_box = ttk.Combobox(root, values=options, textvariable=var_value,
                                 validate="key", validatecommand=cmd)
        if(len(options) > 0):
            if isinstance(var_value, enum_var):
                combo_box.current(var_value.get_enum().value)
            else:
                combo_box.current(0)
        return window._w_create_pair(root, combo_box, text_name, base_x, base_y)

    def _w_create_pair_checkbox(root, text_name, var_value, base_x, base_y):
        cb = tkinter.Checkbutton(root, variable=var_value,
                                 bg=window._color_bg, highlightcolor=window._color_bg, activebackground=window._color_bg)
        return window._w_create_pair(root, cb, text_name, base_x, base_y)

    def _w_create_pair_color_picker(root, text_name, var_value, base_x, base_y, can_be_disabled=True):
        picker = color_picker(root, var_value, window._color_bg, can_be_disabled=can_be_disabled)
        return window._w_create_pair(root, picker, text_name, base_x, base_y)

    def _check_file_name(file_name):
        return file_name is not None and len(file_name) > 0

    def _update_pres_data(self):
        if self.distrib.input_file_paths is not None:
            self._all_set()
        self.edit_source_slide["values"] = self.distrib.slide_names

    def _update_input_data(self):
        if self.distrib.pres is not None:
            self._all_set()

    def _display_not_implemented():
        tkinter.messagebox.showerror("Error", "This function is not yet implemented.")

    def _select_directory(self, directory, label_text=None):
        file_selector = panel_pick_images(window._color_bg, directory)
        file_selector.grab_set()
        self.top.wait_window(file_selector)
        if file_selector.is_ok:
            res = self.distrib.set_input(file_selector.selected_image_files)
            if res is not input_data_load_result.FAILURE:
                if res is input_data_load_result.OK_NON_UNIFORM_DIMENSIONS:
                    tkinter.messagebox.showwarning("Warning", "Some images have different dimensions than others.\n"
                                                              "Will average the dimensions to provide an uniform grid.")
                self.label_dir["text"] = directory if label_text is None else label_text
                self._update_input_data()

    def _cmd_select_file_pptx(self):
        file_name = tkinter.filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx *.pptm")])
        if window._check_file_name(file_name):
            res = self.distrib.set_pres(file_name)
            if res is True:
                self.label_pres["text"] = file_name
                self._update_pres_data()

    def _cmd_select_file_pdf(self):
        file_name = tkinter.filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if window._check_file_name(file_name):
            extractor = pdf_extractor(file_name)
            num_images = extractor.obtain_num_images()
            if num_images > 0:
                var_from = tkinter.IntVar()
                var_to = tkinter.IntVar()
                panel = panel_extract_pdf_options(window._color_bg, num_images, var_from, var_to, self.top)
                panel.grab_set()
                self.top.wait_window(panel)
                if panel.is_ok:
                    extractor.extract(var_from.get() - 1, var_to.get() - 1)
                    self._select_directory(extractor.output_directory, file_name)
            else:
                tkinter.messagebox.showerror("Error", "There are no images in this PDF.")

    def _cmd_select_dir(self):
        file_name = tkinter.filedialog.askdirectory()
        if window._check_file_name(file_name):
            self._select_directory(file_name)

    def _cmd_create(self):
        use_bg = self.edit_bg_color.is_enabled()
        success = self.distrib.create(use_bg)
        if success:
            tkinter.messagebox.showinfo("Success", "Image distribution succeeded.")
            self.edit_source_slide["values"] = self.distrib.slide_names
        else:
            tkinter.messagebox.showerror("Error", "An error has occured.")

    def conditional_exit(self):
        if tkinter.messagebox.askyesno("Are you sure?", "Are you sure you want to exit?"):
            self.top.destroy()

    def loop(self):
        self.top.mainloop()
        # while True:
        # self.top.update_idletasks()
        # self.top.update()
        # time.sleep(0.016)


def main():
    wnd = window()
    wnd.loop()


if __name__ == "__main__":
    main()
